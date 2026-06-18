"""Slide-type renderers that map Story JSON to PowerPoint shapes."""

from __future__ import annotations

from math import cos, sin, tau

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def _text(shape, text, font, size, color, bold=False, align=None):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align:
        p.alignment = align
    p.font.name = font
    p.font.size = size
    p.font.bold = bold
    p.font.color.rgb = color
    return shape


def _title(slide, title, theme, subtitle=None):
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.55), Inches(0.12), Inches(0.72))
    accent.fill.solid(); accent.fill.fore_color.rgb = theme.red; accent.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.82), Inches(0.47), Inches(11.5), Inches(0.8))
    _text(box, title, theme.font_headline, theme.title_size, theme.black, True)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.85), Inches(1.22), Inches(10.9), Inches(0.38))
        _text(sub, subtitle, theme.font_body, theme.subtitle_size, theme.mid_grey)


def _footer(slide, slide_no, theme):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.18), Inches(13.333), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = theme.red; line.line.fill.background()
    box = slide.shapes.add_textbox(Inches(11.7), Inches(7.22), Inches(1.1), Inches(0.2))
    _text(box, f"{slide_no:02d}", theme.font_body, theme.small_size, theme.mid_grey, align=PP_ALIGN.RIGHT)


def _card(slide, x, y, w, h, headline, body, theme, fill=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill or theme.white; shape.line.color.rgb = theme.border_grey
    tf = shape.text_frame; tf.clear(); tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.16); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]; p.text = headline; p.font.name = theme.font_headline; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = theme.red
    p2 = tf.add_paragraph(); p2.text = body; p2.font.name = theme.font_body; p2.font.size = theme.body_size; p2.font.color.rgb = theme.dark_grey
    return shape


def _bullets(slide, items, x, y, w, h, theme):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.name = theme.font_body; p.font.size = theme.body_size; p.font.color.rgb = theme.dark_grey; p.space_after = Pt(8)
    return box


def _center(shape, text, theme, size=15, color=None, bold=True):
    return _text(shape, text, theme.font_headline if bold else theme.font_body, Pt(size), color or theme.dark_grey, bold, PP_ALIGN.CENTER)


def render_slide(slide, slide_json, theme, safe_area, slide_no: int) -> None:
    renderer = RENDERERS.get(slide_json["type"])
    if renderer is None:
        raise ValueError(f"Unsupported slide type: {slide_json['type']}")
    renderer(slide, slide_json, theme, safe_area)
    _footer(slide, slide_no, theme)


def render_cover(slide, data, theme, safe_area):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    band.fill.solid(); band.fill.fore_color.rgb = theme.red; band.line.fill.background()
    label = slide.shapes.add_textbox(Inches(0.75), Inches(0.34), Inches(4.0), Inches(0.3))
    _text(label, "KNAUF | PRESENTATION ENGINE", theme.font_body, Pt(11), theme.white, True)
    title = slide.shapes.add_textbox(Inches(0.9), Inches(1.82), Inches(9.8), Inches(1.05))
    _text(title, data["title"], theme.font_headline, Pt(42), theme.black, True)
    sub = slide.shapes.add_textbox(Inches(0.95), Inches(2.9), Inches(8.5), Inches(0.4))
    _text(sub, data.get("subtitle", ""), theme.font_body, theme.subtitle_size, theme.mid_grey)
    _card(slide, Inches(0.95), Inches(4.2), Inches(7.2), Inches(1.18), "Kernebudskab", "Udvidet scope. Hybrid værdiprofil. Tid til at afstemme rolle, titel og løn.", theme, theme.light_grey)


def render_executive_summary(slide, data, theme, safe_area):
    _title(slide, data["title"], theme)
    for i, card in enumerate(data["cards"][:3]):
        x = Inches(0.95 + i * 4.05)
        _card(slide, x, Inches(2.35), Inches(3.45), Inches(2.35), card["headline"], card["body"], theme, theme.light_grey if i == 1 else theme.white)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.12), Inches(2.05), Inches(0.58), Inches(0.58))
        dot.fill.solid(); dot.fill.fore_color.rgb = theme.red; dot.line.fill.background(); _center(dot, str(i + 1), theme, 16, theme.white)


def render_timeline(slide, data, theme, safe_area):
    _title(slide, data["title"], theme)
    items = data["items"]
    y = Inches(3.38)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.15), y, Inches(11.0), Inches(0.05))
    line.fill.solid(); line.fill.fore_color.rgb = theme.red; line.line.fill.background()
    gap = 10.5 / max(1, len(items) - 1)
    for idx, item in enumerate(items):
        x = Inches(1.0 + idx * gap)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y - Inches(0.18), Inches(0.42), Inches(0.42))
        dot.fill.solid(); dot.fill.fore_color.rgb = theme.red; dot.line.color.rgb = theme.white
        _card(slide, x - Inches(0.55), Inches(2.0 if idx % 2 == 0 else 3.95), Inches(2.15), Inches(1.05), item["label"], item["detail"], theme)


def render_comparison(slide, data, theme, safe_area):
    _title(slide, data["title"], theme)
    _card(slide, Inches(0.95), Inches(2.15), Inches(5.15), Inches(3.15), data["left_title"], "", theme, theme.light_grey)
    _bullets(slide, data["left_items"], Inches(1.35), Inches(3.0), Inches(4.25), Inches(1.55), theme)
    _card(slide, Inches(7.2), Inches(2.15), Inches(5.15), Inches(3.15), data["right_title"], "", theme)
    _bullets(slide, data["right_items"], Inches(7.55), Inches(3.0), Inches(4.25), Inches(1.55), theme)
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.28), Inches(3.35), Inches(0.78), Inches(0.5))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = theme.red; arrow.line.fill.background()


def render_value_cards(slide, data, theme, safe_area):
    _title(slide, data["title"], theme)
    for i, card in enumerate(data["cards"][:4]):
        _card(slide, Inches(0.95 + (i % 2) * 5.95), Inches(2.1 + (i // 2) * 1.72), Inches(5.25), Inches(1.2), card["headline"], card["body"], theme, theme.light_grey if i % 2 == 0 else theme.white)


def render_house_model(slide, data, theme, safe_area):
    _title(slide, data["title"], theme)
    roof = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.0), Inches(1.85), Inches(7.2), Inches(1.35))
    roof.fill.solid(); roof.fill.fore_color.rgb = theme.red; roof.line.color.rgb = theme.white; _center(roof, data["roof"], theme, 18, theme.white)
    for i, level in enumerate([part.strip() for part in data["upper"].split("|")]):
        block = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.15), Inches(3.05 + i * 0.62), Inches(8.9), Inches(0.48))
        block.fill.solid(); block.fill.fore_color.rgb = theme.light_grey if i % 2 == 0 else theme.white; block.line.color.rgb = theme.red; _center(block, level, theme)
    core = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.15), Inches(4.98), Inches(8.9), Inches(0.62))
    core.fill.solid(); core.fill.fore_color.rgb = theme.dark_grey; core.line.color.rgb = theme.white; _center(core, data["core"], theme, 17, theme.white)
    foundation = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.65), Inches(5.82), Inches(9.9), Inches(0.58))
    foundation.fill.solid(); foundation.fill.fore_color.rgb = theme.mid_grey; foundation.line.color.rgb = theme.white; _center(foundation, data["foundation"], theme, 15, theme.white)
    if data.get("side_note"):
        _card(slide, Inches(9.85), Inches(4.1), Inches(2.3), Inches(1.05), "Note", data["side_note"], theme)


def render_case_list(slide, data, theme, safe_area):
    _title(slide, data["title"], theme)
    for i, case in enumerate(data["cases"][:3]):
        _card(slide, Inches(1.0 + i * 3.75), Inches(2.35), Inches(3.25), Inches(2.1), case["headline"], case["body"], theme, theme.light_grey if i == 0 else theme.white)


def render_hub_and_spokes(slide, data, theme, safe_area):
    _title(slide, data["title"], theme)
    cx, cy = Inches(6.25), Inches(3.75)
    hub = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - Inches(0.95), cy - Inches(0.55), Inches(1.9), Inches(1.1))
    hub.fill.solid(); hub.fill.fore_color.rgb = theme.red; hub.line.fill.background(); _center(hub, data["center"], theme, 16, theme.white)
    items = data["items"]
    for i, item in enumerate(items):
        angle = tau * i / len(items) - tau / 4
        x = cx + Inches(3.4 * cos(angle)); y = cy + Inches(1.9 * sin(angle))
        _card(slide, x - Inches(1.15), y - Inches(0.45), Inches(2.3), Inches(0.9), str(i + 1), item, theme, theme.light_grey if i % 2 == 0 else theme.white)


def render_conclusion(slide, data, theme, safe_area):
    _title(slide, data["title"], theme)
    msg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(2.0), Inches(11.15), Inches(1.25))
    msg.fill.solid(); msg.fill.fore_color.rgb = theme.red; msg.line.fill.background(); _center(msg, data["final_message"], theme, 20, theme.white)
    for i, statement in enumerate(data["statements"][:3]):
        _card(slide, Inches(1.0 + i * 3.75), Inches(4.35), Inches(3.25), Inches(1.05), f"Pointe {i + 1}", statement, theme)


RENDERERS = {
    "cover": render_cover,
    "executive_summary": render_executive_summary,
    "timeline": render_timeline,
    "comparison": render_comparison,
    "value_cards": render_value_cards,
    "house_model": render_house_model,
    "case_list": render_case_list,
    "hub_and_spokes": render_hub_and_spokes,
    "conclusion": render_conclusion,
}
