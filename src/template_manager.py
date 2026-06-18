"""Template loading and safety helpers for PowerPoint rendering."""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE_PATH = ROOT / "templates" / "Knauf.pptx"


@dataclass(frozen=True)
class SafeArea:
    left: Any = Inches(0.85)
    top: Any = Inches(0.55)
    width: Any = Inches(11.65)
    height: Any = Inches(6.25)
    footer_top: Any = Inches(7.18)


class TemplateManager:
    """Loads the Knauf template, removes demo slides and exposes rendering boundaries."""

    def __init__(self, template_path: Path = TEMPLATE_PATH):
        self.template_path = template_path
        self.safe_area = SafeArea()

    def load(self):
        prs = Presentation(str(self.template_path)) if self.template_path.exists() else Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        self.remove_sample_slides(prs)
        return prs

    def remove_sample_slides(self, prs) -> None:
        slide_id_list = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public delete API.
        for slide_id in list(slide_id_list):
            rel_id = slide_id.rId
            prs.part.drop_rel(rel_id)
            slide_id_list.remove(slide_id)

    def best_layout(self, prs, slide_type: str):
        for layout in prs.slide_layouts:
            if "blank" in layout.name.lower():
                return layout
        return prs.slide_layouts[-1]

    def detect_footer(self, prs) -> bool:
        return True

    def get_safe_area(self) -> SafeArea:
        return self.safe_area
