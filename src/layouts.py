"""Reusable slide layout primitives built on python-pptx."""

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .theme import (
    BLACK,
    BODY_SIZE,
    BORDER_GREY,
    DARK_GREY,
    FONT_BODY,
    FONT_HEADLINE,
    KNAUF_RED,
    LIGHT_GREY,
    MID_GREY,
    SECTION_SIZE,
    SMALL_SIZE,
    SUBTITLE_SIZE,
    TITLE_SIZE,
    WHITE,
)


def set_wide(prs):
    """Force a 16:9 widescreen canvas; safe for most corporate templates."""
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)


def blank_layout(prs):
    """Return a blank layout, falling back to the last template layout."""
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[-1]


def add_footer(slide, slide_no: int):
    """Add a minimal footer line and page number."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.18), Inches(13.333), Inches(0.03))
    line.fill.solid(); line.fill.fore_color.rgb = KNAUF_RED
    line.line.fill.background()
    box = slide.shapes.add_textbox(Inches(11.7), Inches(7.22), Inches(1.1), Inches(0.2))
    p = box.text_frame.paragraphs[0]
    p.text = f"{slide_no:02d}"
    p.font.name = FONT_BODY; p.font.size = SMALL_SIZE; p.font.color.rgb = MID_GREY
    p.alignment = PP_ALIGN.RIGHT


def add_title(slide, title: str, subtitle: str | None = None):
    """Add a large consulting-style headline and optional subtitle."""
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(0.55), Inches(0.12), Inches(0.72))
    accent.fill.solid(); accent.fill.fore_color.rgb = KNAUF_RED
    accent.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.82), Inches(0.47), Inches(11.5), Inches(0.8))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = FONT_HEADLINE; p.font.size = TITLE_SIZE; p.font.bold = True; p.font.color.rgb = BLACK
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.85), Inches(1.22), Inches(10.9), Inches(0.38))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = FONT_BODY; sp.font.size = SUBTITLE_SIZE; sp.font.color.rgb = MID_GREY


def add_card(slide, x, y, w, h, title, body, fill=WHITE):
    """Draw a rounded card with title and body copy."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER_GREY
    tf = shape.text_frame; tf.clear(); tf.margin_left = Inches(0.18); tf.margin_right = Inches(0.16); tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = title; p.font.name = FONT_HEADLINE; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = KNAUF_RED
    p2 = tf.add_paragraph(); p2.text = body; p2.font.name = FONT_BODY; p2.font.size = BODY_SIZE; p2.font.color.rgb = DARK_GREY
    return shape


def add_bullets(slide, items, x, y, w, h, color=DARK_GREY):
    """Add concise bullet copy."""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.level = 0
        p.font.name = FONT_BODY; p.font.size = BODY_SIZE; p.font.color.rgb = color
        p.space_after = Pt(8)
    return box


def add_section_label(slide, text, x, y, w):
    box = slide.shapes.add_textbox(x, y, w, Inches(0.28))
    p = box.text_frame.paragraphs[0]
    p.text = text.upper(); p.font.name = FONT_BODY; p.font.size = SMALL_SIZE; p.font.bold = True; p.font.color.rgb = KNAUF_RED
    return box
