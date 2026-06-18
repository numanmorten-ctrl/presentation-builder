"""PowerPoint generation service for Executive Deck Studio."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .content_salary import DEFAULT_DECK_STYLE, DEFAULT_PROMPT, DEFAULT_SLIDE_COUNT, OUTPUT_FILENAME, build_salary_story
from .diagrams import add_house_model, add_step_arrow, add_timeline
from .layouts import add_bullets, add_card, add_footer, add_section_label, add_title, blank_layout, set_wide
from .theme import BLACK, DARK_GREY, FONT_BODY, FONT_HEADLINE, KNAUF_RED, LIGHT_GREY, MID_GREY, WHITE

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE_PATH = ROOT / "templates" / "Knauf.pptx"
OUTPUT_DIR = ROOT / "output"
OUTPUT_PATH = OUTPUT_DIR / OUTPUT_FILENAME


def _clear_template_slides(prs):
    """Remove any example slides while retaining template masters/layouts."""
    slide_id_list = prs.slides._sldIdLst  # noqa: SLF001 - python-pptx has no public delete API.
    for slide_id in list(slide_id_list):
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        slide_id_list.remove(slide_id)


def _new_slide(prs, index: int, title: str | None = None, subtitle: str | None = None):
    slide = prs.slides.add_slide(blank_layout(prs))
    if title:
        add_title(slide, title, subtitle)
    add_footer(slide, index)
    return slide


def _add_notes(slide, notes: str):
    if not notes:
        return
    slide.name = notes[:250]
    box = slide.shapes.add_textbox(Inches(-5), Inches(-5), Inches(1), Inches(1))
    box.name = "Speaker notes"
    box.text_frame.text = notes


def _cover(prs, slide_no, data):
    slide = _new_slide(prs, slide_no)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.0))
    band.fill.solid(); band.fill.fore_color.rgb = KNAUF_RED; band.line.fill.background()
    label = slide.shapes.add_textbox(Inches(0.75), Inches(0.34), Inches(4.0), Inches(0.3))
    p = label.text_frame.paragraphs[0]; p.text = data["kicker"]; p.font.name = FONT_BODY; p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = WHITE
    title = slide.shapes.add_textbox(Inches(0.9), Inches(1.82), Inches(9.8), Inches(1.05))
    p = title.text_frame.paragraphs[0]; p.text = data["title"]; p.font.name = FONT_HEADLINE; p.font.size = Pt(42); p.font.bold = True; p.font.color.rgb = BLACK
    sub = slide.shapes.add_textbox(Inches(0.95), Inches(2.9), Inches(8.5), Inches(0.4))
    p = sub.text_frame.paragraphs[0]; p.text = data["subtitle"]; p.font.name = FONT_BODY; p.font.size = Pt(18); p.font.color.rgb = MID_GREY
    add_card(slide, Inches(0.95), Inches(4.2), Inches(7.2), Inches(1.18), "Core message", data["key_message"], fill=LIGHT_GREY)
    add_card(slide, Inches(8.55), Inches(4.2), Inches(3.15), Inches(1.18), "Mode", "Rule-based draft\nAI-ready workflow", fill=WHITE)
    return slide


def _summary(prs, slide_no, data):
    slide = _new_slide(prs, slide_no, data["title"], data.get("subtitle"))
    for i, (head, body) in enumerate(data["points"]):
        x = Inches(0.95 + i * 4.05)
        card = add_card(slide, x, Inches(2.35), Inches(3.45), Inches(2.35), head, body, fill=WHITE if i != 1 else LIGHT_GREY)
        num = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.12), Inches(2.05), Inches(0.58), Inches(0.58))
        num.fill.solid(); num.fill.fore_color.rgb = KNAUF_RED; num.line.fill.background()
        p = num.text_frame.paragraphs[0]; p.text = str(i + 1); p.alignment = PP_ALIGN.CENTER; p.font.name = FONT_HEADLINE; p.font.bold = True; p.font.size = Pt(16); p.font.color.rgb = WHITE
    return slide


def _comparison(prs, slide_no, data):
    slide = _new_slide(prs, slide_no, data["title"], data.get("subtitle"))
    add_card(slide, Inches(0.95), Inches(2.15), Inches(5.15), Inches(3.15), data["left_title"], "", fill=LIGHT_GREY)
    add_bullets(slide, data["left"], Inches(1.35), Inches(3.0), Inches(4.25), Inches(1.55))
    add_card(slide, Inches(7.2), Inches(2.15), Inches(5.15), Inches(3.15), data["right_title"], "", fill=WHITE)
    add_bullets(slide, data["right"], Inches(7.55), Inches(3.0), Inches(4.25), Inches(1.55))
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.28), Inches(3.35), Inches(0.78), Inches(0.5))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = KNAUF_RED; arrow.line.fill.background()
    return slide


def _cards(prs, slide_no, data):
    slide = _new_slide(prs, slide_no, data["title"], data.get("subtitle"))
    for i, (heading, body) in enumerate(data["cards"]):
        x = Inches(0.95 + (i % 2) * 5.95); y = Inches(2.1 + (i // 2) * 1.72)
        add_card(slide, x, y, Inches(5.25), Inches(1.2), heading, body, fill=WHITE if i % 2 else LIGHT_GREY)
    return slide


def _results(prs, slide_no, data):
    slide = _new_slide(prs, slide_no, data["title"], data.get("subtitle"))
    for i, (big, label) in enumerate(data["metrics"]):
        x = Inches(0.85 + i * 3.05)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.35), Inches(2.5), Inches(1.9))
        card.fill.solid(); card.fill.fore_color.rgb = KNAUF_RED if i == 0 else WHITE; card.line.color.rgb = KNAUF_RED
        tf = card.text_frame; tf.clear()
        p = tf.paragraphs[0]; p.text = big; p.alignment = PP_ALIGN.CENTER; p.font.name = FONT_HEADLINE; p.font.size = Pt(25); p.font.bold = True; p.font.color.rgb = WHITE if i == 0 else KNAUF_RED
        p2 = tf.add_paragraph(); p2.text = label; p2.alignment = PP_ALIGN.CENTER; p2.font.name = FONT_BODY; p2.font.size = Pt(15); p2.font.color.rgb = WHITE if i == 0 else DARK_GREY
    add_section_label(slide, "Evidence to attach", Inches(1.0), Inches(5.1), Inches(4.5))
    add_bullets(slide, data["evidence"], Inches(1.0), Inches(5.45), Inches(10.5), Inches(0.8))
    return slide


def _conclusion(prs, slide_no, data):
    slide = _new_slide(prs, slide_no, data["title"])
    msg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(2.0), Inches(11.15), Inches(1.25))
    msg.fill.solid(); msg.fill.fore_color.rgb = KNAUF_RED; msg.line.fill.background()
    p = msg.text_frame.paragraphs[0]; p.text = data["message"]; p.alignment = PP_ALIGN.CENTER; p.font.name = FONT_HEADLINE; p.font.size = Pt(21); p.font.bold = True; p.font.color.rgb = WHITE
    for i, ask in enumerate(data["asks"]):
        add_card(slide, Inches(1.0 + i * 3.75), Inches(4.35), Inches(3.25), Inches(1.05), f"Ask {i+1}", ask, fill=WHITE)
    return slide


def build_presentation(
    prompt: str = DEFAULT_PROMPT,
    deck_style: str = DEFAULT_DECK_STYLE,
    slide_count: int = DEFAULT_SLIDE_COUNT,
    include_notes: bool = True,
    output_path: Path | None = None,
) -> Path:
    """Build the executive salary/development deck and return the generated file path."""
    output_path = output_path or OUTPUT_PATH
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(TEMPLATE_PATH)) if TEMPLATE_PATH.exists() else Presentation()
    _clear_template_slides(prs)
    set_wide(prs)
    story = build_salary_story(prompt, deck_style, slide_count, include_notes)
    builders = {"cover": _cover, "summary": _summary, "comparison": _comparison, "cards": _cards, "results": _results, "conclusion": _conclusion}
    for idx, data in enumerate(story, start=1):
        if data["kind"] == "timeline":
            slide = _new_slide(prs, idx, data["title"], data.get("subtitle")); add_timeline(slide, data["items"])
        elif data["kind"] == "house":
            slide = _new_slide(prs, idx, data["title"], data.get("subtitle")); add_house_model(slide, data["foundation"], data["core"], data["levels"], data["roof"])
        elif data["kind"] == "potential":
            slide = _new_slide(prs, idx, data["title"], data.get("subtitle")); add_step_arrow(slide, data["steps"])
        else:
            slide = builders[data["kind"]](prs, idx, data)
        _add_notes(slide, data.get("notes", ""))
    prs.save(output_path)
    return output_path
