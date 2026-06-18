"""Executive diagrams used by the deck builder."""

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .layouts import add_card
from .theme import BODY_SIZE, DARK_GREY, FONT_BODY, FONT_HEADLINE, KNAUF_RED, LIGHT_GREY, MID_GREY, WHITE


def _center_text(shape, text, size=16, color=DARK_GREY, bold=True):
    tf = shape.text_frame; tf.clear(); tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]; p.text = text; p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_HEADLINE if bold else FONT_BODY; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color


def add_timeline(slide, items):
    y = Inches(3.38)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.15), y, Inches(11.0), Inches(0.05))
    line.fill.solid(); line.fill.fore_color.rgb = KNAUF_RED; line.line.fill.background()
    gap = 10.5 / (len(items) - 1)
    for idx, (year, heading, body) in enumerate(items):
        x = Inches(1.0 + idx * gap)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y - Inches(0.18), Inches(0.42), Inches(0.42))
        dot.fill.solid(); dot.fill.fore_color.rgb = KNAUF_RED; dot.line.color.rgb = WHITE
        box_y = Inches(2.0 if idx % 2 == 0 else 3.95)
        add_card(slide, x - Inches(0.55), box_y, Inches(2.15), Inches(1.05), year, f"{heading}\n{body}", fill=WHITE)


def add_step_arrow(slide, steps):
    for i, step in enumerate(steps):
        x = Inches(0.8 + i * 3.05)
        shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, Inches(3.12), Inches(2.75), Inches(1.05))
        shp.fill.solid(); shp.fill.fore_color.rgb = KNAUF_RED if i == 0 else LIGHT_GREY
        shp.line.color.rgb = WHITE
        _center_text(shp, step, size=14, color=WHITE if i == 0 else DARK_GREY)


def add_house_model(slide, foundation, core, levels, roof):
    """Draw a consulting-style house model for an interdisciplinary capability stack."""
    left = Inches(2.15); width = Inches(8.9)
    roof_shape = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(3.0), Inches(1.85), Inches(7.2), Inches(1.35))
    roof_shape.fill.solid(); roof_shape.fill.fore_color.rgb = KNAUF_RED; roof_shape.line.color.rgb = WHITE
    _center_text(roof_shape, roof, size=18, color=WHITE)

    for i, level in enumerate(levels):
        y = Inches(3.05 + i * 0.62)
        block = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, y, width, Inches(0.48))
        block.fill.solid(); block.fill.fore_color.rgb = WHITE if i % 2 else LIGHT_GREY
        block.line.color.rgb = KNAUF_RED if i == 0 else MID_GREY
        _center_text(block, level, size=15, color=DARK_GREY)

    core_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(4.98), width, Inches(0.62))
    core_block.fill.solid(); core_block.fill.fore_color.rgb = DARK_GREY; core_block.line.color.rgb = WHITE
    _center_text(core_block, core, size=17, color=WHITE)

    foundation_block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.65), Inches(5.82), Inches(9.9), Inches(0.58))
    foundation_block.fill.solid(); foundation_block.fill.fore_color.rgb = MID_GREY; foundation_block.line.color.rgb = WHITE
    _center_text(foundation_block, foundation, size=15, color=WHITE)

    for x in (1.95, 10.95):
        pillar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(3.08), Inches(0.18), Inches(2.72))
        pillar.fill.solid(); pillar.fill.fore_color.rgb = KNAUF_RED; pillar.line.fill.background()
